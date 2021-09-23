from pptx import Presentation
from pptx.util import Inches

ADD_PPTX_NAME = r"/diffExample.pptx"


class PowerPointWorker:
    def __init__(self, way):
        self.__way = way + ADD_PPTX_NAME
        self.__prs = Presentation()

    def insert_info(self, text, image_way="test.jpg"):
        slide = self.__prs.slides.add_slide(self.__prs.slide_layouts[0])  # adding a slide

        title = slide.placeholders[0]  # placeholder for subtitle

        title.text = text  # subtitle

        slide2 = self.__prs.slides.add_slide(self.__prs.slide_layouts[6])

        left = Inches(1)
        top = Inches(0.5)
        slide2.shapes.add_picture(image_way, left, top)

    def save(self):
        self.__prs.save(self.__way)

    def set_way(self, way):
        self.__way = way + ADD_PPTX_NAME
